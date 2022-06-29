#!/usr/bin/env python3

from typing import List, Tuple

import tempfile

from PIL import Image
import pdf2image

from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE

SLIDE_LAYOUT_BLANK = 6

def ComputeImageRect(slide_width: int, slide_height: int,
                     image_width: int, image_height: int
                     ) -> Tuple[int, int, int, int]:
    slide_ratio = slide_width / slide_height
    image_ratio = image_width / image_height
    if slide_ratio > image_ratio:
        # ┌─┬───────┬─┐
        # │ │       │ │
        # │ │ image │ │
        # │ │       │ │
        # └─┴───────┴─┘
        top = 0
        height = slide_height
        width = int(height * image_ratio)
        left = (slide_width - width) / 2
    else:
        # ┌───────┐
        # ├───────┤
        # │       │
        # │ image │
        # │       │
        # ├───────┤
        # └───────┘
        left = 0
        width = slide_width
        height = int(width / image_ratio)
        top = (slide_height - height) / 2
    return (left, top, width, height)

def AppendImageSlide(image: Image.Image, presentation: Presentation) -> None:
    layout = presentation.slide_layouts[SLIDE_LAYOUT_BLANK]
    slide = presentation.slides.add_slide(layout)
    with tempfile.NamedTemporaryFile(suffix='.jpg') as imgfile:
        left, top, width, height = ComputeImageRect(
            slide_width=presentation.slide_width,
            slide_height=presentation.slide_height,
            image_width=image.width,
            image_height=image.height)

        image.save(imgfile.name, "JPEG")

        slide.shapes.add_picture(imgfile.name,
                                 left=left, top=top, width=width, height=height)

def MaxImageSize(images: List[Image.Image]) -> Tuple[int, int]:
    width = max(img.width for img in images)
    height = max(img.height for img in images)
    return width, height

def RenderPptx(filename: str, images: List[Image.Image]) -> None:
    max_width_pt, max_height_pt = MaxImageSize(images)

    presentation = Presentation()
    presentation.slide_width = Pt(max_width_pt)
    presentation.slide_height = Pt(max_height_pt)

    for img in images:
        AppendImageSlide(img, presentation)
    presentation.save(filename)

def ConvertPdfToPptx(pdf_filename: str, pptx_filename: str, dpi: int):
    images = pdf2image.convert_from_path(pdf_filename, dpi=dpi)
    RenderPptx(pptx_filename, images)

if __name__ == "__main__":
    import os
    import sys
    import logging

    for in_file in sys.argv[1:]:
        if not in_file.endswith(".pdf"):
            continue

        filename, extension = os.path.splitext(in_file)
        out_file = filename + ".pptx"

        print(f"{in_file} -> {out_file}")
        try:
            ConvertPdfToPptx(pdf_filename=in_file,
                             pptx_filename=out_file,
                             dpi=120)
        except:
            logging.exception("Failed to convert {in_file}")
